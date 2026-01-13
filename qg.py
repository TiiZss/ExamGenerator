import os
import argparse
import pypdf
import docx
from pptx import Presentation
import subprocess
import time
import sys
import platform

def get_gemini_api_key():
    """Obtiene la clave de API de Google Gemini desde una variable de entorno."""
    # Intentar desde settings.json primero (para web)
    try:
        from examgenerator.utils.settings import get_gemini_api_key as get_key_from_settings
        api_key = get_key_from_settings()
        if api_key:
            return api_key
    except ImportError:
        pass
    
    # Fallback a variable de entorno
    api_key = os.environ.get('GOOGLE_API_KEY')
    if not api_key:
        raise ValueError("Error: La variable de entorno GOOGLE_API_KEY no está configurada. "
                         "Por favor, configúrala con tu clave de API de Google AI.")
    return api_key

def extract_text_from_pdf(file_obj):
    """Extrae texto de un objeto fichero PDF (stream o path)."""
    print(f"📄 Extrayendo texto del PDF...")
    try:
        reader = pypdf.PdfReader(file_obj)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Error al leer el PDF: {e}")
        return None

def extract_text_from_docx(file_obj):
    """Extrae texto de un objeto fichero DOCX (stream o path)."""
    print(f"📄 Extrayendo texto del DOCX...")
    try:
        doc = docx.Document(file_obj)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error al leer el DOCX: {e}")
        return None

def extract_text_from_pptx(file_obj):
    """Extrae texto de un objeto fichero PPTX (stream o path)."""
    print(f"📄 Extrayendo texto del PPTX...")
    try:
        prs = Presentation(file_obj)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error al leer el PPTX: {e}")
        return None

def check_ollama_running(ollama_url):
    """Verifica si el servidor Ollama está corriendo."""
    try:
        import requests
        response = requests.get(f"{ollama_url}/api/tags", timeout=2)
        return response.status_code == 200
    except:
        return False

def start_ollama_server():
    """Intenta iniciar el servidor Ollama automáticamente."""
    print("🚀 Intentando iniciar el servidor Ollama...")
    
    system = platform.system()
    
    try:
        if system == "Windows":
            # En Windows, Ollama se instala como servicio pero también se puede ejecutar directamente
            # Intentar iniciar el servicio primero
            try:
                subprocess.run(["net", "start", "Ollama"], 
                             capture_output=True, text=True, timeout=5)
            except:
                pass
            
            # Si no funciona como servicio, intentar ejecutar directamente
            try:
                subprocess.Popen(["ollama", "serve"], 
                               stdout=subprocess.DEVNULL, 
                               stderr=subprocess.DEVNULL,
                               creationflags=subprocess.CREATE_NEW_CONSOLE if system == "Windows" else 0)
            except FileNotFoundError:
                print("❌ Error: Ollama no está instalado o no está en PATH")
                print("   Descarga Ollama desde: https://ollama.ai")
                return False
                
        else:  # Linux/macOS
            try:
                # Intentar iniciar en segundo plano
                subprocess.Popen(["ollama", "serve"], 
                               stdout=subprocess.DEVNULL, 
                               stderr=subprocess.DEVNULL)
            except FileNotFoundError:
                print("❌ Error: Ollama no está instalado o no está en PATH")
                print("   Instala Ollama desde: https://ollama.ai")
                return False
        
        print("⏳ Esperando a que Ollama inicie...")
        # Esperar hasta 30 segundos para que el servidor inicie
        for i in range(30):
            time.sleep(1)
            if check_ollama_running("http://localhost:11434"):
                print("✅ Servidor Ollama iniciado correctamente")
                return True
            if i % 5 == 0 and i > 0:
                print(f"   Esperando... ({i}s)")
        
        print("⚠️  Ollama no respondió en 30 segundos")
        return False
        
    except Exception as e:
        print(f"❌ Error al intentar iniciar Ollama: {e}")
        return False

def ensure_ollama_running(ollama_url, interactive=True):
    """Asegura que Ollama esté corriendo, iniciándolo si es necesario.
    
    Args:
        ollama_url: URL del servidor Ollama
        interactive: Si es True, pregunta al usuario. Si es False, solo verifica sin interacción.
    """
    if check_ollama_running(ollama_url):
        print(f"✅ Servidor Ollama ya está corriendo en {ollama_url}")
        return True
    
    print(f"⚠️  Servidor Ollama no detectado en {ollama_url}")
    
    # En modo no-interactivo, no intentar iniciar
    if not interactive:
        print("❌ Ollama no está disponible (modo no-interactivo)")
        print("   Para usar Ollama, inicia el servicio manualmente o usa el perfil 'ollama' de Docker")
        return False
    
    # Solo en modo interactivo preguntar
    try:
        import sys
        # Verificar si stdin está disponible
        if not sys.stdin.isatty():
            print("❌ No hay terminal interactiva disponible")
            print("   Inicia Ollama manualmente con: ollama serve")
            return False
            
        response = input("¿Quieres que intente iniciar Ollama automáticamente? (s/n): ").strip().lower()
        
        if response in ['s', 'y', 'si', 'yes']:
            return start_ollama_server()
        else:
            print("❌ No se puede continuar sin Ollama")
            print("   Inicia Ollama manualmente con: ollama serve")
            return False
    except (EOFError, OSError):
        print("❌ No se puede leer entrada del usuario")
        print("   Inicia Ollama manualmente con: ollama serve")
        return False

def generate_questions_with_gemini(text_content, num_questions=10, language="español", model_name="gemini-1.5-flash"):
    """Genera preguntas usando Google Gemini con chunking automático para grandes volúmenes."""
    try:
        import google.generativeai as genai
    except ImportError:
        print("❌ Error: Necesitas instalar google-generativeai para usar Gemini:")
        print("pip install google-generativeai")
        return None
    
    if not text_content or text_content.strip() == "":
        print("El texto extraído está vacío. No se pueden generar preguntas.")
        return None

    # Configurar chunking: Gemini puede manejar más preguntas por llamada
    MAX_QUESTIONS_PER_CALL = 15
    
    if num_questions <= MAX_QUESTIONS_PER_CALL:
        # Caso simple: una sola llamada
        return _generate_gemini_chunk(text_content, num_questions, language, model_name)
    
    # Caso complejo: múltiples llamadas
    print(f"\n🧠 Generando {num_questions} preguntas con Google Gemini ({model_name})...")
    print(f"📦 Dividiendo en chunks de {MAX_QUESTIONS_PER_CALL} preguntas...")
    
    all_questions = []
    chunks = (num_questions + MAX_QUESTIONS_PER_CALL - 1) // MAX_QUESTIONS_PER_CALL  # Redondeo hacia arriba
    
    for chunk_idx in range(chunks):
        questions_in_chunk = min(MAX_QUESTIONS_PER_CALL, num_questions - chunk_idx * MAX_QUESTIONS_PER_CALL)
        print(f"\n🔄 Chunk {chunk_idx + 1}/{chunks}: Generando {questions_in_chunk} preguntas...")
        
        chunk_result = _generate_gemini_chunk(text_content, questions_in_chunk, language, model_name)
        
        if chunk_result:
            all_questions.append(chunk_result)
        else:
            print(f"⚠️ Error en chunk {chunk_idx + 1}, continuando con las generadas...")
    
    if not all_questions:
        return None
    
    # Combinar resultados
    combined = "\n\n".join(all_questions)
    print(f"\n✅ Total: {num_questions} preguntas generadas exitosamente")
    return combined


def _generate_gemini_chunk(text_content, num_questions, language, model_name):
    """Genera un chunk de preguntas con Gemini (función interna)."""
    try:
        import google.generativeai as genai
        genai.configure(api_key=get_gemini_api_key())
        model = genai.GenerativeModel(model_name)
        
        prompt = f"""
        Basándote en el siguiente texto, genera exactamente {num_questions} preguntas de opción múltiple en formato AIKEN.
        
        FORMATO REQUERIDO (ejemplo):
        ¿Pregunta aquí?
        A) Opción 1
        B) Opción 2
        C) Opción 3
        D) Opción 4
        ANSWER: B)
        
        REGLAS IMPORTANTES:
        - Cada pregunta debe tener exactamente 4 opciones (A, B, C, D)
        - Las opciones deben usar el formato: A), B), C), D)
        - La respuesta correcta debe indicarse con: ANSWER: X) donde X es A, B, C o D
        - Las preguntas deben ser claras y evaluar la comprensión del contenido
        - Genera las preguntas en {language}
        - Deja una línea en blanco entre cada pregunta

        Texto base:
        ---
        {text_content}
        ---
        """
        
        response = model.generate_content(prompt)
        return response.text
        
    except Exception as e:
        print(f"\n❌ Error al contactar con la API de Gemini: {e}")
        return None

def generate_questions_with_ollama(text_content, num_questions=10, language="español", 
                                   model_name="llama2", ollama_url="http://localhost:11434",
                                   interactive=True):
    """Genera preguntas usando Ollama (IA local) con chunking automático para grandes volúmenes.
    
    Args:
        text_content: Texto para generar preguntas
        num_questions: Número de preguntas a generar
        language: Idioma de las preguntas
        model_name: Modelo de Ollama a usar
        ollama_url: URL del servidor Ollama
        interactive: Si es False, no intenta iniciar Ollama interactivamente
    """
    try:
        import requests
    except ImportError:
        print("❌ Error: Necesitas instalar requests para usar Ollama:")
        print("pip install requests")
        return None
    
    if not text_content or text_content.strip() == "":
        print("El texto extraído está vacío. No se pueden generar preguntas.")
        return None

    # Verificar y asegurar que Ollama esté corriendo
    if not ensure_ollama_running(ollama_url, interactive=interactive):
        return None
    
    # Configurar chunking: Con 2000 tokens, asumimos ~10 preguntas por llamada
    MAX_QUESTIONS_PER_CALL = 10
    
    if num_questions <= MAX_QUESTIONS_PER_CALL:
        # Caso simple: una sola llamada
        print(f"\n🤖 Enviando texto a Ollama local ({model_name})...")
        print(f"📡 URL de Ollama: {ollama_url}")
        return _generate_ollama_chunk(text_content, num_questions, language, model_name, ollama_url)
    
    # Caso complejo: múltiples llamadas
    print(f"\n🤖 Generando {num_questions} preguntas con Ollama ({model_name})...")
    print(f"📦 Dividiendo en chunks de {MAX_QUESTIONS_PER_CALL} preguntas...")
    print(f"📡 URL de Ollama: {ollama_url}")
    
    all_questions = []
    chunks = (num_questions + MAX_QUESTIONS_PER_CALL - 1) // MAX_QUESTIONS_PER_CALL  # Redondeo hacia arriba
    
    for chunk_idx in range(chunks):
        questions_in_chunk = min(MAX_QUESTIONS_PER_CALL, num_questions - chunk_idx * MAX_QUESTIONS_PER_CALL)
        print(f"\n🔄 Chunk {chunk_idx + 1}/{chunks}: Generando {questions_in_chunk} preguntas...")
        
        chunk_result = _generate_ollama_chunk(text_content, questions_in_chunk, language, model_name, ollama_url)
        
        if chunk_result:
            all_questions.append(chunk_result)
        else:
            print(f"⚠️ Error en chunk {chunk_idx + 1}, continuando con las generadas...")
    
    if not all_questions:
        return None
    
    # Combinar resultados
    combined = "\n\n".join(all_questions)
    print(f"\n✅ Total: {num_questions} preguntas generadas exitosamente")
    return combined


def _generate_ollama_chunk(text_content, num_questions, language, model_name, ollama_url):
    """Genera un chunk de preguntas con Ollama (función interna)."""
    try:
        import requests
        
        prompt = f"""
        TAREA: Genera EXACTAMENTE {num_questions} preguntas de opción múltiple en formato AIKEN sobre el texto proporcionado.
        
        FORMATO OBLIGATORIO para CADA pregunta:
        ¿Pregunta aquí?
        A) Opción 1
        B) Opción 2
        C) Opción 3
        D) Opción 4
        ANSWER: B)
        
        (línea en blanco)
        
        REGLAS ESTRICTAS:
        1. DEBES generar {num_questions} preguntas completas
        2. Cada pregunta DEBE tener exactamente 4 opciones (A, B, C, D)
        3. Cada pregunta DEBE tener su línea ANSWER: X)
        4. Deja UNA línea en blanco entre preguntas
        5. Idioma: {language}
        6. NO incluyas numeración (1., 2., etc.) antes de las preguntas
        7. COMPLETA TODAS las {num_questions} preguntas antes de terminar

        Texto base:
        ---
        {text_content}
        ---
        
        IMPORTANTE: Genera las {num_questions} preguntas completas ahora.
        """
        
        # Hacer la petición a la API de Ollama
        response = requests.post(
            f"{ollama_url}/api/generate",
            json={
                "model": model_name,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "num_predict": 2000,  # Tokens suficientes para múltiples preguntas
                    "temperature": 0.7,
                    "top_p": 0.9
                }
            },
            timeout=300  # 5 minutos de timeout
        )
        
        if response.status_code == 200:
            result = response.json()
            return result.get('response', '')
        else:
            print(f"❌ Error en la respuesta de Ollama: {response.status_code}")
            print(f"Respuesta: {response.text}")
            return None
            
    except Exception as e:
        print(f"\n❌ Error al contactar con Ollama: {e}")
        print(f"\n❌ Error al contactar con Ollama: {e}")
        return None


def main():
    """Función principal para orquestar el proceso."""
    parser = argparse.ArgumentParser(
        description="Genera preguntas sobre un fichero PDF, DOCX o PPTX usando IA (Gemini u Ollama).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Ejemplos de uso:
  # Con Google Gemini (por defecto)
  python qg.py documento.pdf --num_preguntas 10
  
  # Con Ollama (IA local)
  python qg.py documento.pdf --motor ollama --modelo llama2 --num_preguntas 10
  
  # Con Gemini especificando modelo
  python qg.py documento.pdf --motor gemini --modelo gemini-1.5-pro --num_preguntas 15
  
  # Con Ollama personalizado
  python qg.py documento.pdf --motor ollama --modelo mistral --ollama_url http://localhost:11434
        """
    )
    parser.add_argument("fichero", help="Ruta al fichero PDF, DOCX o PPTX.")
    parser.add_argument("--num_preguntas", type=int, default=10, 
                       help="Número de preguntas a generar (por defecto: 10).")
    parser.add_argument("--idioma", type=str, default="español", 
                       help="Idioma para las preguntas (por defecto: español).")
    parser.add_argument("--motor", type=str, choices=["gemini", "ollama"], default="gemini",
                       help="Motor de IA a usar: gemini (Google) u ollama (local) (por defecto: gemini).")
    parser.add_argument("--modelo", type=str, default=None,
                       help="Modelo específico a usar. Gemini: gemini-1.5-flash (default), gemini-1.5-pro. Ollama: llama2 (default), mistral, codellama, etc.")
    parser.add_argument("--ollama_url", type=str, default="http://localhost:11434",
                       help="URL del servidor Ollama (por defecto: http://localhost:11434).")
    
    args = parser.parse_args()
    
    file_path = args.fichero
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if not os.path.exists(file_path):
        print(f"❌ Error: El fichero '{file_path}' no existe.")
        return

    # Extraer texto según el tipo de archivo
    extracted_text = ""
    if file_ext == '.pdf':
        extracted_text = extract_text_from_pdf(file_path)
    elif file_ext == '.docx':
        extracted_text = extract_text_from_docx(file_path)
    elif file_ext == '.pptx':
        extracted_text = extract_text_from_pptx(file_path)
    else:
        print(f"❌ Error: Tipo de fichero '{file_ext}' no soportado. Usa PDF, DOCX o PPTX.")
        return

    if not extracted_text:
        print("❌ No se pudo extraer texto del fichero.")
        return

    print(f"✅ Texto extraído exitosamente ({len(extracted_text)} caracteres)")

    # Determinar el modelo por defecto según el motor
    if args.modelo is None:
        model_name = "gemini-1.5-flash" if args.motor == "gemini" else "llama2"
    else:
        model_name = args.modelo

    # Generar preguntas según el motor seleccionado
    questions = None
    if args.motor == "gemini":
        questions = generate_questions_with_gemini(
            extracted_text, 
            args.num_preguntas, 
            args.idioma,
            model_name
        )
    elif args.motor == "ollama":
        questions = generate_questions_with_ollama(
            extracted_text, 
            args.num_preguntas, 
            args.idioma,
            model_name,
            args.ollama_url
        )

    if questions:
        print("\n✅ ¡Aquí tienes las preguntas generadas! ✅")
        print("-" * 60)
        print(questions)
        print("-" * 60)
        print(f"\n📊 Motor usado: {args.motor.upper()} | Modelo: {model_name}")
    else:
        print("\n❌ No se pudieron generar preguntas.")

if __name__ == "__main__":
    main()